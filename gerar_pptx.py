"""
Gerador automatico de PPTX - Dashboard Financeiro FPGolfe
Usa o template original como base e gera slides com dados atualizados
"""
import os
import copy
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import numpy as np

# Cores padrão do template
AZUL_ESCURO = RGBColor(0x33, 0x33, 0x33)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
PRETO = RGBColor(0x00, 0x00, 0x00)
CINZA = RGBColor(0x66, 0x66, 0x66)
CINZA_CLARO = RGBColor(0xF2, 0xF2, 0xF2)
VERDE = RGBColor(0x2E, 0x7D, 0x32)
VERMELHO = RGBColor(0xC6, 0x28, 0x28)
AZUL = RGBColor(0x0D, 0x47, 0xA1)
HEADER_BG = RGBColor(0x33, 0x33, 0x33)
ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)

# Dimensões do slide (widescreen 13.33 x 7.5 inches)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def fmt_brl(val, mil=False):
    """Formata valor em R$"""
    if pd.isna(val) or val == 0:
        return "R$ 0"
    if mil:
        if abs(val) >= 1_000_000:
            return f"R$ {val/1_000_000:,.1f}M".replace(",", "X").replace(".", ",").replace("X", ".")
        elif abs(val) >= 1_000:
            return f"R$ {val/1_000:,.0f}K".replace(",", "X").replace(".", ",").replace("X", ".")
    return f"R$ {val:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")


def add_header_bar(slide, section_num, title_text):
    """Adiciona barra de cabeçalho escura no topo do slide (igual ao template)"""
    # Barra escura no topo
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()
    
    # Número da seção
    txBox = slide.shapes.add_textbox(Cm(0.8), Cm(0.3), Cm(1.5), Cm(1))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"0{section_num}"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Public Sans"
    run.font.color.rgb = BRANCO
    
    # Título
    txBox2 = slide.shapes.add_textbox(Cm(0.8), Cm(1.2), Cm(20), Cm(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = title_text
    run2.font.size = Pt(18)
    run2.font.bold = True
    run2.font.name = "Public Sans"
    run2.font.color.rgb = BRANCO
    
    # Linha separadora
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(2.5), SLIDE_W, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL
    line.line.fill.background()


def add_subtitle(slide, text, left=Cm(0.8), top=Cm(3), width=Cm(30)):
    """Adiciona subtítulo abaixo do header"""
    txBox = slide.shapes.add_textbox(left, top, width, Cm(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.name = "Public Sans"
    run.font.color.rgb = CINZA
    return txBox


def add_table(slide, data, col_widths, left, top, width, row_height=Cm(0.7),
              header_bg=HEADER_BG, header_fg=BRANCO, font_size=Pt(9)):
    """Adiciona tabela formatada ao slide"""
    rows = len(data)
    cols = len(data[0]) if data else 0
    if rows == 0 or cols == 0:
        return None
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Emu(int(row_height) * rows))
    table = table_shape.table
    
    # Ajustar largura das colunas
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_val) if cell_val is not None else ""
            
            # Formatação do texto
            for para in cell.text_frame.paragraphs:
                para.font.size = font_size
                para.font.name = "Public Sans"
                
                if r_idx == 0:  # Header
                    para.font.bold = True
                    para.font.color.rgb = header_fg
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.color.rgb = PRETO
                    # Alinhar números à direita
                    if c_idx > 0 and isinstance(cell_val, str) and ("R$" in cell_val or "%" in cell_val or cell_val.replace(".", "").replace(",", "").replace("-", "").isdigit()):
                        para.alignment = PP_ALIGN.RIGHT
                    elif c_idx == 0:
                        para.alignment = PP_ALIGN.LEFT
                    else:
                        para.alignment = PP_ALIGN.RIGHT
            
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Cores de fundo
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRANCO
    
    return table_shape


def gerar_apresentacao(df_razao, df_orcamento, grupo_desc, conta_grupo, conta_desc,
                       template_path=None, trimestre="1T26", periodo_meses=[1, 2, 3]):
    """
    Gera a apresentação PPTX completa.
    
    Args:
        df_razao: DataFrame com razão contábil processado
        df_orcamento: DataFrame com orçamento
        grupo_desc: dict grupo -> descrição
        conta_grupo: dict conta -> grupo
        conta_desc: dict conta -> descrição
        template_path: caminho do PPTX template original (não usado - recriamos tudo)
        trimestre: string do trimestre (ex: "1T26")
        periodo_meses: lista de meses inclusos
    """
    # Calcular dados
    meses_nome = {1: "Jan", 2: "Fev", 3: "Mar", 4: "Abr", 5: "Mai", 6: "Jun",
                  7: "Jul", 8: "Ago", 9: "Set", 10: "Out", 11: "Nov", 12: "Dez"}
    
    df_periodo = df_razao[df_razao["MesNum"].isin(periodo_meses)].copy()
    
    # Totais
    rec_total = df_periodo[df_periodo["Tipo"] == "Receita"]["Valor"].sum()
    desp_total = df_periodo[df_periodo["Tipo"] == "Despesa"]["Valor"].sum()
    resultado = rec_total - desp_total
    
    # Orçamento do período
    df_orc_periodo = df_orcamento[df_orcamento["MesNum"].isin(periodo_meses)].copy()
    rec_orc = df_orc_periodo[df_orc_periodo["Tipo"] == "Receita"]["ValorOrc"].sum()
    desp_orc = df_orc_periodo[df_orc_periodo["Tipo"] == "Despesa"]["ValorOrc"].sum()
    resultado_orc = rec_orc - desp_orc
    
    # Criar PPTX novo usando template como referência de dimensões
    prs_new = Presentation()
    prs_new.slide_width = SLIDE_W
    prs_new.slide_height = SLIDE_H
    blank = prs_new.slide_layouts[6]  # Em branco
    
    # ===== SLIDE 1: CAPA =====
    slide_capa = prs_new.slides.add_slide(blank)
    
    # Imagem de fundo (campo de golfe)
    capa_img_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "capa_fundo.jpg")
    if os.path.exists(capa_img_path):
        slide_capa.shapes.add_picture(capa_img_path, 0, 0, SLIDE_W, SLIDE_H)
    else:
        # Fallback cor sólida
        bg = slide_capa.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x1B)
    
    # Overlay escuro semi-transparente (retângulo preto)
    overlay = slide_capa.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    overlay.line.fill.background()
    # Aplicar transparência via XML
    from lxml import etree
    spPr = overlay._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if spPr is not None:
        srgb = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '50000')  # 50% transparência
    
    # Logo do golfista (canto direito)
    logo_img_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "capa_logo.png")
    if os.path.exists(logo_img_path):
        slide_capa.shapes.add_picture(logo_img_path, Cm(27), Cm(0.5), Cm(5), Cm(8.5))
    
    # Título "FPGolfe"
    txBox = slide_capa.shapes.add_textbox(Cm(2), Cm(3), Cm(20), Cm(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "FPGolfe"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.name = "Public Sans"
    run.font.color.rgb = BRANCO
    
    # Subtítulo
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    periodo_texto = f"Resultado {trimestre}"
    run2.text = periodo_texto
    run2.font.size = Pt(28)
    run2.font.name = "Public Sans"
    run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    
    # ===== SLIDE 2: ÍNDICE (idêntico ao original) =====
    slide_idx = prs_new.slides.add_slide(blank)
    
    # Golfista preto grande no lado esquerdo
    golfista_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "golfista_preto.png")
    if os.path.exists(golfista_path):
        slide_idx.shapes.add_picture(golfista_path, Cm(0), Cm(0.5), Cm(8), Cm(17))
    
    # Título "RESULTADOS 1T26" - lado direito, acima
    txTitle = slide_idx.shapes.add_textbox(Cm(10), Cm(0.8), Cm(20), Cm(1.5))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"RESULTADOS {trimestre}"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Public Sans"
    run.font.color.rgb = PRETO
    
    # Itens do índice - posicionados à direita do golfista
    items = [
        ("0 1", "BP / DRE"),
        ("0 2", "DRE por Centro de Custo"),
        ("0 3", "Apreciações e pontos relevantes"),
    ]
    
    y_start = Cm(3.2)
    for i, (num, desc) in enumerate(items):
        # Número grande cinza
        txNum = slide_idx.shapes.add_textbox(Cm(10), y_start + Cm(i * 3.2), Cm(5), Cm(1.5))
        tf_num = txNum.text_frame
        p_num = tf_num.paragraphs[0]
        run_num = p_num.add_run()
        run_num.text = num
        run_num.font.size = Pt(24)
        run_num.font.bold = False
        run_num.font.name = "Public Sans"
        run_num.font.color.rgb = CINZA
        
        # Descrição abaixo do número
        txDesc = slide_idx.shapes.add_textbox(Cm(10), y_start + Cm(i * 3.2 + 1.2), Cm(18), Cm(1))
        tf_desc = txDesc.text_frame
        p_desc = tf_desc.paragraphs[0]
        run_desc = p_desc.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(14)
        run_desc.font.name = "Public Sans"
        run_desc.font.color.rgb = PRETO
        
        # Linha separadora abaixo
        line = slide_idx.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10), y_start + Cm(i * 3.2 + 2.2), Cm(18), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        line.line.fill.background()
    
    # Logos na parte inferior (GOLFE Federação Paulista + FPG golf center)
    logo_fpg_vermelho = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo_fpgolfe_vermelho.png")
    logo_fpg_golf_center = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo_fpg_golf_center.png")
    if os.path.exists(logo_fpg_vermelho):
        slide_idx.shapes.add_picture(logo_fpg_vermelho, Cm(12.5), Cm(14), Cm(3.5), Cm(3.5))
    if os.path.exists(logo_fpg_golf_center):
        slide_idx.shapes.add_picture(logo_fpg_golf_center, Cm(17.5), Cm(14.8), Cm(4), Cm(1.8))
    
    # ===== SLIDE 3: DRE =====
    slide_dre = prs_new.slides.add_slide(blank)
    add_header_bar(slide_dre, 1, f"BP / DRE | {trimestre}")
    add_subtitle(slide_dre, f"A DRE apresenta um resultado {'superavitário' if resultado > 0 else 'deficitário'} de {fmt_brl(abs(resultado), mil=True)}.")
    
    # Montar DRE
    # Receitas por grupo
    rec_grupos = df_periodo[df_periodo["Tipo"] == "Receita"].groupby("DescGrupo")["Valor"].sum().sort_values(ascending=False)
    desp_grupos = df_periodo[df_periodo["Tipo"] == "Despesa"].groupby("DescGrupo")["Valor"].sum().sort_values(ascending=False)
    
    dre_data = [["DEMONSTRAÇÃO DO RESULTADO", "PERÍODO", "ACUMULADO"]]
    
    dre_data.append(["RECEITAS TOTAIS", fmt_brl(rec_total), fmt_brl(rec_total)])
    for grp, val in rec_grupos.items():
        dre_data.append([f"   {grp}", fmt_brl(val), fmt_brl(val)])
    
    dre_data.append(["", "", ""])
    dre_data.append(["DESPESAS TOTAIS", fmt_brl(desp_total), fmt_brl(desp_total)])
    for grp, val in desp_grupos.items():
        dre_data.append([f"   {grp}", fmt_brl(val), fmt_brl(val)])
    
    dre_data.append(["", "", ""])
    dre_data.append(["SUPERÁVIT/DÉFICIT", fmt_brl(resultado), fmt_brl(resultado)])
    
    col_widths = [Cm(14), Cm(6), Cm(6)]
    add_table(slide_dre, dre_data, col_widths, Cm(1), Cm(4), Cm(26), font_size=Pt(9))
    
    # ===== SLIDE 4: DRE MENSAL =====
    slide_dre_mes = prs_new.slides.add_slide(blank)
    add_header_bar(slide_dre_mes, 1, f"DRE MENSAL | {trimestre}")
    add_subtitle(slide_dre_mes, "Evolução mensal de Receitas, Despesas e Resultado.")
    
    # Tabela mensal
    header = [""] + [meses_nome[m] for m in periodo_meses] + ["ACUMULADO"]
    
    mensal_data = [header]
    
    # Receita por mês
    rec_mensal = ["RECEITAS"]
    rec_acum = 0
    for m in periodo_meses:
        val = df_periodo[(df_periodo["MesNum"] == m) & (df_periodo["Tipo"] == "Receita")]["Valor"].sum()
        rec_mensal.append(fmt_brl(val))
        rec_acum += val
    rec_mensal.append(fmt_brl(rec_acum))
    mensal_data.append(rec_mensal)
    
    # Despesa por mês
    desp_mensal = ["DESPESAS"]
    desp_acum = 0
    for m in periodo_meses:
        val = df_periodo[(df_periodo["MesNum"] == m) & (df_periodo["Tipo"] == "Despesa")]["Valor"].sum()
        desp_mensal.append(fmt_brl(val))
        desp_acum += val
    desp_mensal.append(fmt_brl(desp_acum))
    mensal_data.append(desp_mensal)
    
    # Resultado por mês
    res_mensal = ["RESULTADO"]
    for m in periodo_meses:
        r = df_periodo[(df_periodo["MesNum"] == m) & (df_periodo["Tipo"] == "Receita")]["Valor"].sum()
        d = df_periodo[(df_periodo["MesNum"] == m) & (df_periodo["Tipo"] == "Despesa")]["Valor"].sum()
        res_mensal.append(fmt_brl(r - d))
    res_mensal.append(fmt_brl(resultado))
    mensal_data.append(res_mensal)
    
    n_cols = len(header)
    col_w = [Cm(6)] + [Cm(5)] * (n_cols - 1)
    total_w = sum(c for c in col_w)
    add_table(slide_dre_mes, mensal_data, col_w, Cm(1), Cm(4), Emu(int(total_w)), font_size=Pt(10))
    
    # ===== SLIDES: DRE POR CENTRO DE CUSTO - FPG / CE / JUV =====
    # Usa a coluna CentroCusto que já vem mapeada do dashboard (De/Para aplicado)
    if "CentroCusto" in df_razao.columns:
        centros = [
            ("FPG", "FPGOLFE", df_periodo[df_periodo["CentroCusto"] == "FPG"].copy()),
            ("CE", "CENTRO ESPORTIVO", df_periodo[df_periodo["CentroCusto"] == "CE"].copy()),
            ("JUV", "JUVENIS", df_periodo[df_periodo["CentroCusto"] == "JUV"].copy()),
        ]
    else:
        # Fallback pela lógica antiga (CR=100001)
        COL_CR = "CR"
        cr_col = COL_CR if COL_CR in df_razao.columns else None
        if cr_col:
            centros = [
                ("FPG", "FPGOLFE", df_periodo[df_periodo[cr_col] != 100001].copy()),
                ("CE", "CENTRO ESPORTIVO", df_periodo[df_periodo[cr_col] == 100001].copy()),
            ]
        else:
            centros = [("FPG", "FPGOLFE", df_periodo.copy())]
    
    for cc_code, cc_nome, df_cc in centros:
        if df_cc.empty:
            continue
        
        slide_cc = prs_new.slides.add_slide(blank)
        add_header_bar(slide_cc, 2, f"DRE POR CENTRO DE CUSTO | {trimestre}")
        add_subtitle(slide_cc, f"Seguem os resultados apartados pelos Centros de Custos: {cc_nome}")
        
        # Título do centro de custo
        txCC = slide_cc.shapes.add_textbox(Cm(1), Cm(4), Cm(30), Cm(1.2))
        tf_cc = txCC.text_frame
        p_cc = tf_cc.paragraphs[0]
        p_cc.alignment = PP_ALIGN.CENTER
        run_cc = p_cc.add_run()
        run_cc.text = cc_nome
        run_cc.font.size = Pt(18)
        run_cc.font.bold = True
        run_cc.font.name = "Public Sans"
        run_cc.font.color.rgb = PRETO
        
        p_sub = tf_cc.add_paragraph()
        p_sub.alignment = PP_ALIGN.CENTER
        run_sub = p_sub.add_run()
        run_sub.text = f"Orçado x Realizado - YTD {periodo_meses[-1]:02d}/2026"
        run_sub.font.size = Pt(11)
        run_sub.font.name = "Public Sans"
        run_sub.font.color.rgb = CINZA
        
        # Calcular valores do centro de custo
        cc_rec = df_cc[df_cc["Tipo"] == "Receita"]["Valor"].sum()
        cc_desp = df_cc[df_cc["Tipo"] == "Despesa"]["Valor"].sum()
        cc_resultado = cc_rec - cc_desp
        
        # Orçamento proporcional
        prop_rec = cc_rec / rec_total if rec_total > 0 else 1.0 / len(centros)
        prop_desp = cc_desp / desp_total if desp_total > 0 else 1.0 / len(centros)
        cc_rec_orc = rec_orc * prop_rec
        cc_desp_orc = desp_orc * prop_desp
        cc_res_orc = cc_rec_orc - cc_desp_orc
        
        # Tabela comparativa
        cc_header = ["TIPO", "ORÇADO", "REALIZADO", "VARIAÇÃO", "DESCRIÇÃO"]
        cc_data = [cc_header]
        
        cc_rec_grupos = df_cc[df_cc["Tipo"] == "Receita"].groupby("DescGrupo")["Valor"].sum().sort_values(ascending=False)
        
        cc_data.append(["TOTAL DAS RECEITAS", fmt_brl(cc_rec_orc), fmt_brl(cc_rec), fmt_brl(cc_rec - cc_rec_orc), "Receitas sociais e torneios"])
        
        for grp, val in cc_rec_grupos.items():
            cc_data.append([f"   {grp}", "", fmt_brl(val), "", ""])
        
        cc_data.append(["", "", "", "", ""])
        
        cc_desp_grupos = df_cc[df_cc["Tipo"] == "Despesa"].groupby("DescGrupo")["Valor"].sum().sort_values(ascending=False)
        
        cc_data.append(["TOTAL DAS DESPESAS", fmt_brl(cc_desp_orc), fmt_brl(cc_desp), fmt_brl(cc_desp - cc_desp_orc), ""])
        
        for grp, val in list(cc_desp_grupos.items())[:8]:
            cc_data.append([f"   {grp}", "", fmt_brl(val), "", ""])
        
        cc_data.append(["", "", "", "", ""])
        cc_data.append(["RESULTADO CONSOLIDADO", fmt_brl(cc_res_orc), fmt_brl(cc_resultado), fmt_brl(cc_resultado - cc_res_orc), ""])
        
        col_w_cc = [Cm(8), Cm(4.5), Cm(4.5), Cm(4.5), Cm(6)]
        add_table(slide_cc, cc_data, col_w_cc, Cm(1), Cm(6), Cm(27.5), font_size=Pt(8), row_height=Cm(0.6))
    
    # ===== SLIDE: ORÇADO VS REALIZADO =====
    slide_oxr = prs_new.slides.add_slide(blank)
    add_header_bar(slide_oxr, 2, f"ORÇADO versus REALIZADO | {trimestre}")
    add_subtitle(slide_oxr, "Seguem os resultados consolidados e comparativo com o Orçado.")
    
    oxr_header = ["", "ORÇADO", "REALIZADO", "VARIAÇÃO", "VAR %"]
    oxr_data = [oxr_header]
    
    var_rec = rec_total - rec_orc
    var_desp = desp_total - desp_orc
    var_res = resultado - resultado_orc
    
    pct_rec = (var_rec / rec_orc * 100) if rec_orc != 0 else 0
    pct_desp = (var_desp / desp_orc * 100) if desp_orc != 0 else 0
    pct_res = (var_res / resultado_orc * 100) if resultado_orc != 0 else 0
    
    oxr_data.append(["RECEITAS", fmt_brl(rec_orc), fmt_brl(rec_total), fmt_brl(var_rec), f"{pct_rec:+.1f}%"])
    oxr_data.append(["DESPESAS", fmt_brl(desp_orc), fmt_brl(desp_total), fmt_brl(var_desp), f"{pct_desp:+.1f}%"])
    oxr_data.append(["RESULTADO", fmt_brl(resultado_orc), fmt_brl(resultado), fmt_brl(var_res), f"{pct_res:+.1f}%"])
    
    col_w_oxr = [Cm(5), Cm(5.5), Cm(5.5), Cm(5.5), Cm(3.5)]
    add_table(slide_oxr, oxr_data, col_w_oxr, Cm(1), Cm(4), Cm(25), font_size=Pt(11))
    
    # ===== SLIDE 6: BUDGET & FORECAST =====
    slide_bf = prs_new.slides.add_slide(blank)
    add_header_bar(slide_bf, 2, f"BUDGET & FORECAST | YTD {periodo_meses[-1]:02d}/2026")
    add_subtitle(slide_bf, "Abaixo a projeção de resultado para 31/12/2026")
    
    # Projeção anual baseada no realizado até agora
    meses_realizados = len(periodo_meses)
    rec_anual_proj = rec_total / meses_realizados * 12
    desp_anual_proj = desp_total / meses_realizados * 12
    res_anual_proj = rec_anual_proj - desp_anual_proj
    
    # Orçamento anual total
    rec_orc_anual = df_orcamento[df_orcamento["Tipo"] == "Receita"]["ValorOrc"].sum()
    desp_orc_anual = df_orcamento[df_orcamento["Tipo"] == "Despesa"]["ValorOrc"].sum()
    res_orc_anual = rec_orc_anual - desp_orc_anual
    
    bf_header = ["", "ORÇAMENTO 2026", f"REALIZADO {trimestre}", "FORECAST 2026", "VAR vs ORC"]
    bf_data = [bf_header]
    bf_data.append(["RECEITAS", fmt_brl(rec_orc_anual), fmt_brl(rec_total), fmt_brl(rec_anual_proj), fmt_brl(rec_anual_proj - rec_orc_anual)])
    bf_data.append(["DESPESAS", fmt_brl(desp_orc_anual), fmt_brl(desp_total), fmt_brl(desp_anual_proj), fmt_brl(desp_anual_proj - desp_orc_anual)])
    bf_data.append(["RESULTADO", fmt_brl(res_orc_anual), fmt_brl(resultado), fmt_brl(res_anual_proj), fmt_brl(res_anual_proj - res_orc_anual)])
    
    col_w_bf = [Cm(5), Cm(5), Cm(5), Cm(5), Cm(5)]
    add_table(slide_bf, bf_data, col_w_bf, Cm(1), Cm(4), Cm(25), font_size=Pt(10))
    
    # Nota de rodapé
    txNote = slide_bf.shapes.add_textbox(Cm(1), Cm(7.5), Cm(30), Cm(1))
    tf = txNote.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "*Tal projeção é passível de ajuste e discussão do forecast."
    run.font.size = Pt(9)
    run.font.italics = True
    run.font.color.rgb = CINZA
    run.font.name = "Public Sans"
    
    # ===== SLIDE 7: RECEITAS DETALHADAS =====
    slide_rec = prs_new.slides.add_slide(blank)
    add_header_bar(slide_rec, 1, f"RECEITAS DETALHADAS | {trimestre}")
    add_subtitle(slide_rec, "Composição das receitas por grupo de conta.")
    
    rec_detail = [["GRUPO", "VALOR", "% DO TOTAL"]]
    for grp, val in rec_grupos.items():
        pct = val / rec_total * 100 if rec_total > 0 else 0
        rec_detail.append([grp, fmt_brl(val), f"{pct:.1f}%"])
    rec_detail.append(["TOTAL RECEITAS", fmt_brl(rec_total), "100%"])
    
    col_w_rec = [Cm(12), Cm(6), Cm(4)]
    add_table(slide_rec, rec_detail, col_w_rec, Cm(1), Cm(4), Cm(22), font_size=Pt(10))
    
    # ===== SLIDE 8: DESPESAS DETALHADAS =====
    slide_desp = prs_new.slides.add_slide(blank)
    add_header_bar(slide_desp, 1, f"DESPESAS DETALHADAS | {trimestre}")
    add_subtitle(slide_desp, "Composição das despesas por grupo de conta.")
    
    desp_detail = [["GRUPO", "VALOR", "% DO TOTAL"]]
    for grp, val in desp_grupos.items():
        pct = val / desp_total * 100 if desp_total > 0 else 0
        desp_detail.append([grp, fmt_brl(val), f"{pct:.1f}%"])
    desp_detail.append(["TOTAL DESPESAS", fmt_brl(desp_total), "100%"])
    
    col_w_desp = [Cm(12), Cm(6), Cm(4)]
    add_table(slide_desp, desp_detail, col_w_desp, Cm(1), Cm(4), Cm(22), font_size=Pt(10))
    
    # ===== SLIDE 9: GESTÃO ESTRATÉGICA (placeholder) =====
    slide_gestao = prs_new.slides.add_slide(blank)
    add_header_bar(slide_gestao, 3, "GESTÃO ESTRATÉGICA E FINANCEIRA DO GOLFE PAULISTA")
    
    # Texto principal com dados reais
    margem_pct = resultado / rec_total * 100 if rec_total > 0 else 0
    var_vs_orc = resultado - resultado_orc
    pct_anual = resultado / res_orc_anual * 100 if res_orc_anual > 0 else 0
    
    texto = (f"O {trimestre} 2026 encerrou com superávit consolidado de {fmt_brl(resultado, mil=True)}, "
             f"{'superando' if var_vs_orc > 0 else 'abaixo d'}o orçado em {fmt_brl(abs(var_vs_orc), mil=True)} "
             f"e representando {pct_anual:.0f}% do resultado anual previsto já no primeiro trimestre.")
    
    txBox = slide_gestao.shapes.add_textbox(Cm(8), Cm(3.5), Cm(22), Cm(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(14)
    run.font.name = "Public Sans"
    run.font.color.rgb = PRETO
    
    # Destaques
    txDest = slide_gestao.shapes.add_textbox(Cm(1), Cm(3.5), Cm(7), Cm(5))
    tf2 = txDest.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Números que\ncontam histórias"
    run2.font.size = Pt(24)
    run2.font.bold = True
    run2.font.name = "Public Sans"
    run2.font.color.rgb = PRETO
    
    bullets = ["Disciplina na Gestão", "Governança forte", "Confiança na entidade"]
    for b in bullets:
        p3 = tf2.add_paragraph()
        p3.space_before = Pt(8)
        run3 = p3.add_run()
        run3.text = f"→ {b}"
        run3.font.size = Pt(12)
        run3.font.name = "Public Sans"
        run3.font.color.rgb = CINZA
    
    # ===== SLIDE FINAL: ENCERRAMENTO =====
    slide_final = prs_new.slides.add_slide(blank)
    
    # Usar a imagem original do slide final (logo vermelho GOLFE Federação Paulista)
    logo_final_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo_final_slide.png")
    if os.path.exists(logo_final_path):
        # Fundo preto
        bg_rect = slide_final.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = PRETO
        bg_rect.line.fill.background()
        # Logo centralizado
        logo_w = Cm(22)
        logo_h = Cm(12.2)
        logo_left = Emu(int((SLIDE_W - logo_w) / 2))
        logo_top = Emu(int((SLIDE_H - logo_h) / 2))
        slide_final.shapes.add_picture(logo_final_path, logo_left, logo_top, logo_w, logo_h)
    else:
        # Fallback texto
        bg = slide_final.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PRETO
        
        txBox = slide_final.shapes.add_textbox(Cm(3), Cm(2), Cm(20), Cm(3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "GOLFE"
        run.font.size = Pt(60)
        run.font.bold = True
        run.font.name = "Public Sans"
        run.font.color.rgb = BRANCO
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = "Federação Paulista"
        run2.font.size = Pt(24)
        run2.font.name = "Public Sans"
        run2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    
    return prs_new
